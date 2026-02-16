
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation(output_path):
    prs = Presentation()

    # Helper function to add a slide with title and content
    def add_slide(title_text, content_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.text = content_points[0]
        for point in content_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Bank-Grade Reconciliation AI"
    subtitle.text = "Architecture Overview & Agentic Workflow\nPowered by Oracle 26AI"

    # 2. The Reconciliation Challenge
    add_slide("The Reconciliation Challenge", [
        "Complex 1:N and M:N Matching: One payment covering multiple invoices or vice versa.",
        "Remittance Ambiguity: Unstructured bank text with missing or typo-ridden references.",
        "High Volume & Velocity: Processing thousands of daily transactions at scale.",
        "Fragmented Source Data: Bridging ERP invoices and Bank payments."
    ])

    # 3. System Architecture
    add_slide("System Architecture", [
        "Frontend: Streamlit-based Dashboard for human-in-the-loop validation.",
        "Orchestration Layer: FastAPI backend managing the agentic lifecycle.",
        "Persistence & AI Layer: Oracle 26AI acting as both a relational store and vector engine.",
        "Communication: RESTful APIs connecting services."
    ])

    # 4. Agentic Orchestration Model
    add_slide("Agentic Orchestration Model", [
        "Linear Workflow Execution: Intake -> Matching -> Exception -> Compliance -> Audit.",
        "State Management: Persistent 'ReconciliationState' tracking across the lifecycle.",
        "Intelligent Routing: Handover between specialized agents based on confidence scores.",
        "Autonomous Reasoning: Agents making decisions within defined guardrails."
    ])

    # 5. Specialized Agents
    add_slide("Specialized Agents", [
        "Intake Agent: Fetches structured data from Oracle 26AI tables.",
        "Matching Agent: Executes hybrid matching (Deterministic + Semantic).",
        "Exception Agent: Classifies and analyzes unmatched transactions.",
        "Compliance Agent: Enforces bank-grade guardrails and risk checks.",
        "Audit Agent: Records an immutable trail of all AI decisions."
    ])

    # 6. Hybrid Matching Engine
    add_slide("Hybrid Matching Engine", [
        "Deterministic Step: Exact matching on PO numbers, IDs, and amounts.",
        "Semantic Step: Utilizing LLM-driven vector embeddings for text proximity.",
        "Oracle Vector Search: Offloading heavy semantic similarity to the database layer.",
        "Confidence Scoring: Ranking matches for automated approval vs. manual review."
    ])

    # 7. Oracle 26AI Data Layer
    add_slide("Oracle 26AI Data Layer", [
        "Unified Storage: Handling JSON remittance data alongside relational tables.",
        "Integrated Vector Engine: Native 'VECTOR' data type for high-performance retrieval.",
        "Tenant Isolation: Robust customer_id partitioning across all schemas.",
        "Reliability: ACID compliance for financial transaction integrity."
    ])

    # 8. Human-in-the-Loop
    add_slide("Human-in-the-Loop", [
        "Transparency: Visualizing AI reasoning and confidence for every decision.",
        "Manual Overrides: Direct interface for human operators to fix outliers.",
        "Persistence: All manual decisions are fed back into the audit trail.",
        "Exception Viewer: Dedicated view for high-risk or complex mismatches."
    ])

    # 9. Security & Auditability
    add_slide("Security & Auditability", [
        "Immutable Audit Trail: Capturing every agent action, prompt hash, and version.",
        "Compliance Guardrails: Blocking prohibited actions (e.g., cross-tenant matching).",
        "Encryption: Handling sensitive financial data securely with .env management.",
        "Traceability: Mapping every match back to its specific source and reasoning."
    ])

    # 10. Conclusion
    add_slide("Conclusion", [
        "Precision: Drastic reduction in manual overhead via 80%+ auto-reconciliation.",
        "Scalability: Built on high-performance Oracle & FastAPI foundations.",
        "Future-Ready: Seamlessly integrates new LLMs and matching patterns.",
        "Operational Efficiency: Real-time visibility into financial reconciliation health."
    ])

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    output_file = "d:/Antigravity/Reconcilation/architecture_ppt.pptx"
    create_presentation(output_file)
